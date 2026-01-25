"""Export service for SMSI documents - Multi-format support with PSSIG styling."""
import io
import os
import tempfile
from datetime import datetime
from typing import Dict, Any, Optional, List, BinaryIO
from pathlib import Path
import uuid

from loguru import logger

from app.infrastructure.database.smsi_models import GeneratedDocument, DocumentType


# PSSIG Color schemes for document types
DOCUMENT_COLORS = {
    "DIRECTIVE": {"primary": "#1e3a5f", "secondary": "#2563eb", "accent": "#3b82f6", "name": "Directive Stratégique"},
    "POLICY": {"primary": "#14532d", "secondary": "#16a34a", "accent": "#22c55e", "name": "Politique"},
    "PROCEDURE": {"primary": "#7c2d12", "secondary": "#ea580c", "accent": "#f97316", "name": "Procédure"},
    "REGISTER": {"primary": "#581c87", "secondary": "#9333ea", "accent": "#a855f7", "name": "Registre"},
    "ANNEX": {"primary": "#1e3a5f", "secondary": "#0891b2", "accent": "#06b6d4", "name": "Annexe"},
    "CHECKLIST": {"primary": "#0f766e", "secondary": "#14b8a6", "accent": "#2dd4bf", "name": "Checklist"},
    "REPORT": {"primary": "#b91c1c", "secondary": "#dc2626", "accent": "#ef4444", "name": "Rapport"},
    "MATRIX": {"primary": "#4338ca", "secondary": "#6366f1", "accent": "#818cf8", "name": "Matrice"},
    "TEMPLATE": {"primary": "#7c3aed", "secondary": "#8b5cf6", "accent": "#a78bfa", "name": "Template"},
    "SCHEMA": {"primary": "#0369a1", "secondary": "#0284c7", "accent": "#0ea5e9", "name": "Schéma"},
    "DEFAULT": {"primary": "#374151", "secondary": "#6b7280", "accent": "#9ca3af", "name": "Document"},
}


class ExportService:
    """Service for exporting documents to multiple formats with PSSIG styling."""

    SUPPORTED_FORMATS = ["md", "html", "docx", "pdf", "xlsx", "csv", "pptx"]

    async def export_document(
        self,
        document: GeneratedDocument,
        format: str,
        output_path: Optional[str] = None,
        organization_name: str = "Organisation"
    ) -> Dict[str, Any]:
        """Export a document to the specified format with professional PSSIG styling."""
        if format not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {format}. Supported: {self.SUPPORTED_FORMATS}")

        # Methods that support organization name
        org_methods = {
            "html": lambda d: self._export_html(d, organization_name),
            "docx": lambda d: self._export_docx(d, organization_name),
            "pdf": lambda d: self._export_pdf(d, organization_name),
        }

        # Methods that don't need organization name
        simple_methods = {
            "md": self._export_markdown,
            "xlsx": self._export_xlsx,
            "csv": self._export_csv,
            "pptx": self._export_pptx,
        }

        try:
            if format in org_methods:
                content, filename = await org_methods[format](document)
            else:
                content, filename = await simple_methods[format](document)

            if output_path:
                with open(output_path, 'wb') as f:
                    f.write(content)
                return {
                    "success": True,
                    "format": format,
                    "filename": filename,
                    "path": output_path,
                    "size": len(content)
                }

            return {
                "success": True,
                "format": format,
                "filename": filename,
                "content": content,
                "size": len(content)
            }

        except Exception as e:
            logger.error(f"Export error ({format}): {str(e)}")
            return {
                "success": False,
                "format": format,
                "error": str(e)
            }

    async def _export_markdown(self, document: GeneratedDocument) -> tuple[bytes, str]:
        """Export to Markdown format."""
        content = document.content_markdown or ""

        # Add metadata header
        header = f"""---
title: {document.name}
code: {document.code}
version: {document.version}
type: {document.document_type.value}
status: {document.status.value}
generated: {document.created_at.isoformat()}
---

"""
        full_content = header + content
        filename = f"{document.code}_v{document.version}.md"

        return full_content.encode('utf-8'), filename

    async def _export_html(self, document: GeneratedDocument, organization_name: str = "Organisation") -> tuple[bytes, str]:
        """Export to HTML format with PSSIG professional styling."""
        html_content = document.content_html or self._markdown_to_html(document.content_markdown)

        # Get colors for document type
        doc_type = document.document_type.value if document.document_type else "DEFAULT"
        colors = DOCUMENT_COLORS.get(doc_type, DOCUMENT_COLORS["DEFAULT"])
        type_name = colors["name"]

        # Wrap in HTML document with professional PSSIG CSS
        full_html = f"""<!DOCTYPE html>
<html lang="fr">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{document.name}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Segoe UI', 'Roboto', 'Helvetica Neue', Arial, sans-serif;
            line-height: 1.7;
            color: #1f2937;
            background-color: #ffffff;
        }}

        /* Cover Page / Header */
        .document-header {{
            background: linear-gradient(135deg, {colors["primary"]} 0%, {colors["secondary"]} 100%);
            color: white;
            padding: 60px 50px;
            margin-bottom: 40px;
            position: relative;
            overflow: hidden;
        }}
        .document-header::before {{
            content: '';
            position: absolute;
            top: -50%;
            right: -10%;
            width: 400px;
            height: 400px;
            background: rgba(255,255,255,0.1);
            border-radius: 50%;
        }}
        .organization-name {{
            font-size: 42px;
            font-weight: 700;
            margin-bottom: 10px;
            letter-spacing: -1px;
        }}
        .organization-underline {{
            width: 80px;
            height: 4px;
            background: rgba(255,255,255,0.5);
            border-radius: 2px;
            margin-bottom: 30px;
        }}
        .doc-type-badge {{
            display: inline-block;
            background: rgba(255,255,255,0.2);
            padding: 8px 20px;
            border-radius: 8px;
            font-size: 14px;
            font-weight: 600;
            margin-bottom: 20px;
        }}
        .document-title {{
            font-size: 28px;
            font-weight: 600;
            margin-bottom: 40px;
            max-width: 700px;
            line-height: 1.3;
        }}
        .metadata-grid {{
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 15px;
            margin-bottom: 25px;
        }}
        .metadata-item {{
            background: rgba(255,255,255,0.1);
            padding: 15px;
            border-radius: 10px;
        }}
        .metadata-label {{
            font-size: 11px;
            text-transform: uppercase;
            letter-spacing: 1px;
            opacity: 0.7;
            margin-bottom: 5px;
        }}
        .metadata-value {{
            font-size: 16px;
            font-weight: 600;
        }}
        .classification-badge {{
            display: inline-block;
            background: rgba(234, 179, 8, 0.2);
            color: #fef3c7;
            border: 1px solid rgba(234, 179, 8, 0.3);
            padding: 5px 15px;
            border-radius: 20px;
            font-size: 13px;
            margin-right: 10px;
        }}

        /* Document Body */
        .document-body {{
            max-width: 850px;
            margin: 0 auto;
            padding: 0 50px 50px;
        }}
        h1 {{
            font-size: 32px;
            font-weight: 700;
            margin-top: 50px;
            margin-bottom: 25px;
            padding-bottom: 15px;
            border-bottom: 3px solid {colors["secondary"]};
            color: #111827;
        }}
        h2 {{
            font-size: 24px;
            font-weight: 600;
            margin-top: 40px;
            margin-bottom: 20px;
            color: {colors["primary"]};
            display: flex;
            align-items: center;
            gap: 12px;
        }}
        h2::before {{
            content: '';
            display: inline-block;
            width: 4px;
            height: 28px;
            background: {colors["secondary"]};
            border-radius: 2px;
        }}
        h3 {{
            font-size: 20px;
            font-weight: 600;
            margin-top: 30px;
            margin-bottom: 15px;
            color: {colors["secondary"]};
        }}
        h4 {{
            font-size: 16px;
            font-weight: 600;
            margin-top: 25px;
            margin-bottom: 10px;
            color: #374151;
        }}
        p {{
            margin-bottom: 16px;
            color: #4b5563;
        }}

        /* Lists */
        ul, ol {{
            margin: 20px 0;
            padding-left: 0;
        }}
        ul {{ list-style: none; }}
        ul li {{
            position: relative;
            padding-left: 25px;
            margin-bottom: 10px;
            color: #4b5563;
        }}
        ul li::before {{
            content: '';
            position: absolute;
            left: 0;
            top: 10px;
            width: 8px;
            height: 8px;
            background: {colors["accent"]};
            border-radius: 50%;
        }}
        ol {{
            padding-left: 25px;
        }}
        ol li {{
            margin-bottom: 10px;
            color: #4b5563;
        }}

        /* Tables */
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 25px 0;
            border-radius: 10px;
            overflow: hidden;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}
        th {{
            background: {colors["primary"]};
            color: white;
            font-weight: 600;
            padding: 15px;
            text-align: left;
            font-size: 14px;
        }}
        td {{
            padding: 14px 15px;
            border-bottom: 1px solid #e5e7eb;
            color: #4b5563;
            font-size: 14px;
        }}
        tr:last-child td {{
            border-bottom: none;
        }}
        tr:nth-child(even) {{
            background: #f9fafb;
        }}
        tr:hover {{
            background: #f3f4f6;
        }}

        /* Blockquotes */
        blockquote {{
            margin: 25px 0;
            padding: 20px 25px;
            background: {colors["secondary"]}10;
            border-left: 4px solid {colors["secondary"]};
            border-radius: 0 10px 10px 0;
            font-style: italic;
            color: #4b5563;
        }}

        /* Code */
        code {{
            background: {colors["secondary"]}15;
            color: {colors["primary"]};
            padding: 3px 8px;
            border-radius: 4px;
            font-family: 'Consolas', 'Monaco', monospace;
            font-size: 0.9em;
        }}
        pre {{
            background: #1f2937;
            color: #e5e7eb;
            padding: 20px;
            border-radius: 10px;
            overflow-x: auto;
            margin: 20px 0;
        }}
        pre code {{
            background: none;
            color: inherit;
            padding: 0;
        }}

        /* Horizontal rule */
        hr {{
            border: none;
            height: 1px;
            background: linear-gradient(to right, transparent, #d1d5db, transparent);
            margin: 40px 0;
        }}

        /* Footer */
        .document-footer {{
            max-width: 850px;
            margin: 0 auto;
            padding: 30px 50px;
            border-top: 1px solid {colors["secondary"]}30;
            display: flex;
            justify-content: space-between;
            color: #6b7280;
            font-size: 13px;
        }}
        .footer-org {{
            color: {colors["secondary"]};
            font-weight: 600;
        }}

        /* Print styles */
        @media print {{
            body {{ background: white; }}
            .document-header {{
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
            }}
            th {{
                -webkit-print-color-adjust: exact !important;
                print-color-adjust: exact !important;
            }}
            .document-body {{ max-width: 100%; padding: 0 30px; }}
        }}
    </style>
</head>
<body>
    <div class="document-header">
        <div class="organization-name">{organization_name}</div>
        <div class="organization-underline"></div>
        <div class="doc-type-badge">{type_name}</div>
        <div class="document-title">{document.name}</div>
        <div class="metadata-grid">
            <div class="metadata-item">
                <div class="metadata-label">Code Document</div>
                <div class="metadata-value">{document.code}</div>
            </div>
            <div class="metadata-item">
                <div class="metadata-label">Version</div>
                <div class="metadata-value">{document.version}</div>
            </div>
            <div class="metadata-item">
                <div class="metadata-label">Date</div>
                <div class="metadata-value">{document.created_at.strftime('%d/%m/%Y')}</div>
            </div>
            <div class="metadata-item">
                <div class="metadata-label">Statut</div>
                <div class="metadata-value">{document.status.value.title()}</div>
            </div>
        </div>
        <span class="classification-badge">Classification: Interne</span>
    </div>

    <div class="document-body">
        {html_content}
    </div>

    <div class="document-footer">
        <div>
            <div class="footer-org">{organization_name}</div>
            <div>Document de référence SMSI</div>
        </div>
        <div style="text-align: right;">
            <div>{document.code} - Version {document.version}</div>
            <div>Généré le {document.created_at.strftime('%d/%m/%Y')}</div>
        </div>
    </div>
</body>
</html>"""

        filename = f"{document.code}_v{document.version}.html"
        return full_html.encode('utf-8'), filename

    async def _export_docx(self, document: GeneratedDocument, organization_name: str = "Organisation") -> tuple[bytes, str]:
        """Export to DOCX format with PSSIG professional styling."""
        from docx import Document as DocxDocument
        from docx.shared import Inches, Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = DocxDocument()

        # Get colors for document type
        doc_type = document.document_type.value if document.document_type else "DEFAULT"
        colors = DOCUMENT_COLORS.get(doc_type, DOCUMENT_COLORS["DEFAULT"])
        type_name = colors["name"]

        # Convert hex color to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

        primary_rgb = hex_to_rgb(colors["primary"])
        secondary_rgb = hex_to_rgb(colors["secondary"])

        # Set up document margins
        sections = doc.sections
        for section in sections:
            section.top_margin = Cm(2)
            section.bottom_margin = Cm(2)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(2.5)

        # Set up styles
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)
        style.paragraph_format.line_spacing = 1.5

        # Customize Heading styles
        for i in range(1, 5):
            heading_style = doc.styles[f'Heading {i}']
            heading_style.font.color.rgb = RGBColor(*primary_rgb)
            heading_style.font.name = 'Calibri'
            if i == 1:
                heading_style.font.size = Pt(24)
                heading_style.font.bold = True
            elif i == 2:
                heading_style.font.size = Pt(18)
                heading_style.font.bold = True
                heading_style.font.color.rgb = RGBColor(*secondary_rgb)
            elif i == 3:
                heading_style.font.size = Pt(14)
                heading_style.font.bold = True

        # =====================================================
        # COVER PAGE
        # =====================================================

        # Organization name (large, colored)
        org_para = doc.add_paragraph()
        org_run = org_para.add_run(organization_name)
        org_run.font.size = Pt(36)
        org_run.font.bold = True
        org_run.font.color.rgb = RGBColor(*primary_rgb)
        org_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        org_para.space_after = Pt(10)

        # Decorative line
        line_para = doc.add_paragraph()
        line_run = line_para.add_run("─" * 30)
        line_run.font.color.rgb = RGBColor(*secondary_rgb)
        line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        line_para.space_after = Pt(40)

        # Document type badge
        type_para = doc.add_paragraph()
        type_run = type_para.add_run(f"[ {type_name.upper()} ]")
        type_run.font.size = Pt(12)
        type_run.font.bold = True
        type_run.font.color.rgb = RGBColor(*secondary_rgb)
        type_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        type_para.space_after = Pt(20)

        # Document title
        title_para = doc.add_paragraph()
        title_run = title_para.add_run(document.name)
        title_run.font.size = Pt(26)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(31, 41, 55)  # Gray-800
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_para.space_after = Pt(50)

        # Metadata table
        metadata_table = doc.add_table(rows=4, cols=2)
        metadata_table.alignment = WD_TABLE_ALIGNMENT.CENTER

        metadata = [
            ("Code Document", document.code),
            ("Version", document.version),
            ("Date", document.created_at.strftime('%d/%m/%Y')),
            ("Statut", document.status.value.title())
        ]

        for i, (label, value) in enumerate(metadata):
            # Label cell
            label_cell = metadata_table.rows[i].cells[0]
            label_cell.text = label
            label_para = label_cell.paragraphs[0]
            label_para.runs[0].font.bold = True
            label_para.runs[0].font.color.rgb = RGBColor(107, 114, 128)  # Gray-500

            # Value cell
            value_cell = metadata_table.rows[i].cells[1]
            value_cell.text = str(value)
            value_para = value_cell.paragraphs[0]
            value_para.runs[0].font.bold = True
            value_para.runs[0].font.color.rgb = RGBColor(*primary_rgb)

        doc.add_paragraph()
        doc.add_paragraph()

        # Classification badge
        class_para = doc.add_paragraph()
        class_run = class_para.add_run("Classification: Interne")
        class_run.font.size = Pt(10)
        class_run.font.italic = True
        class_run.font.color.rgb = RGBColor(202, 138, 4)  # Yellow-600
        class_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Page break
        doc.add_page_break()

        # =====================================================
        # DOCUMENT CONTENT
        # =====================================================

        # Add header with metadata
        header = doc.sections[0].header
        header_para = header.paragraphs[0]
        header_run = header_para.add_run(f"{document.code} | Version {document.version} | {organization_name}")
        header_run.font.size = Pt(9)
        header_run.font.color.rgb = RGBColor(107, 114, 128)
        header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

        # Parse and add content with professional styling
        self._parse_markdown_to_docx_styled(doc, document.content_markdown or "", colors)

        # Add footer
        footer = doc.sections[0].footer
        footer_para = footer.paragraphs[0]
        footer_run = footer_para.add_run(f"{organization_name} - Document de référence SMSI - {datetime.now().strftime('%d/%m/%Y')}")
        footer_run.font.size = Pt(9)
        footer_run.font.color.rgb = RGBColor(107, 114, 128)
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        filename = f"{document.code}_v{document.version}.docx"
        return buffer.getvalue(), filename

    def _parse_markdown_to_docx_styled(self, doc, markdown_content: str, colors: dict):
        """Parse Markdown content and add to DOCX document with PSSIG styling."""
        from docx.shared import Pt, RGBColor

        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

        primary_rgb = hex_to_rgb(colors["primary"])
        secondary_rgb = hex_to_rgb(colors["secondary"])

        lines = markdown_content.split('\n')
        in_table = False
        table_data = []

        for line in lines:
            stripped = line.strip()

            # Skip YAML frontmatter
            if stripped == '---':
                continue

            # Headers with colored styling
            if stripped.startswith('# '):
                if in_table:
                    self._add_table_to_docx_styled(doc, table_data, colors)
                    in_table = False
                    table_data = []
                heading = doc.add_heading(stripped[2:], level=1)
                for run in heading.runs:
                    run.font.color.rgb = RGBColor(*primary_rgb)
            elif stripped.startswith('## '):
                heading = doc.add_heading(stripped[3:], level=2)
                for run in heading.runs:
                    run.font.color.rgb = RGBColor(*secondary_rgb)
            elif stripped.startswith('### '):
                heading = doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith('#### '):
                heading = doc.add_heading(stripped[5:], level=4)

            # Tables
            elif stripped.startswith('|'):
                if not in_table:
                    in_table = True
                if not stripped.replace('|', '').replace('-', '').replace(' ', '').replace(':', ''):
                    continue  # Skip separator line
                row = [cell.strip() for cell in stripped.split('|')[1:-1]]
                if row:
                    table_data.append(row)

            # Lists
            elif stripped.startswith('- ') or stripped.startswith('* '):
                text = stripped[2:]
                para = doc.add_paragraph(text, style='List Bullet')
            elif stripped.startswith('1. ') or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == '.'):
                text = stripped.split('. ', 1)[1] if '. ' in stripped else stripped
                para = doc.add_paragraph(text, style='List Number')

            # Regular paragraphs
            elif stripped:
                if in_table:
                    self._add_table_to_docx_styled(doc, table_data, colors)
                    in_table = False
                    table_data = []
                para = doc.add_paragraph()
                self._add_formatted_text(para, stripped)

            # Empty line
            else:
                if in_table:
                    self._add_table_to_docx_styled(doc, table_data, colors)
                    in_table = False
                    table_data = []

        # Handle remaining table
        if in_table and table_data:
            self._add_table_to_docx_styled(doc, table_data, colors)

    def _add_table_to_docx_styled(self, doc, table_data: List[List[str]], colors: dict):
        """Add a professionally styled table to the DOCX document."""
        from docx.shared import Pt, RGBColor
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        if not table_data:
            return

        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

        primary_rgb = hex_to_rgb(colors["primary"])

        rows = len(table_data)
        cols = max(len(row) for row in table_data)

        table = doc.add_table(rows=rows, cols=cols)
        table.style = 'Table Grid'

        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                if j < cols:
                    cell = table.rows[i].cells[j]
                    cell.text = cell_text

                    # Style header row with primary color
                    if i == 0:
                        # Set background color
                        shading_elm = OxmlElement('w:shd')
                        shading_elm.set(qn('w:fill'), colors["primary"].lstrip('#'))
                        cell._tc.get_or_add_tcPr().append(shading_elm)

                        # Set text to white and bold
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(255, 255, 255)

        doc.add_paragraph()

    def _parse_markdown_to_docx(self, doc, markdown_content: str):
        """Parse Markdown content and add to DOCX document."""
        from docx.shared import Pt

        lines = markdown_content.split('\n')
        in_table = False
        table_data = []
        in_list = False
        list_items = []

        for line in lines:
            stripped = line.strip()

            # Headers
            if stripped.startswith('# '):
                if in_table:
                    self._add_table_to_docx(doc, table_data)
                    in_table = False
                    table_data = []
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith('## '):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith('### '):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith('#### '):
                doc.add_heading(stripped[5:], level=4)

            # Tables
            elif stripped.startswith('|'):
                if not in_table:
                    in_table = True
                if not stripped.replace('|', '').replace('-', '').replace(' ', ''):
                    continue  # Skip separator line
                row = [cell.strip() for cell in stripped.split('|')[1:-1]]
                if row:
                    table_data.append(row)

            # Lists
            elif stripped.startswith('- ') or stripped.startswith('* '):
                text = stripped[2:]
                para = doc.add_paragraph(text, style='List Bullet')
            elif stripped.startswith('1. ') or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == '.'):
                text = stripped.split('. ', 1)[1] if '. ' in stripped else stripped
                para = doc.add_paragraph(text, style='List Number')

            # Regular paragraphs
            elif stripped:
                if in_table:
                    self._add_table_to_docx(doc, table_data)
                    in_table = False
                    table_data = []
                # Handle bold and italic
                para = doc.add_paragraph()
                self._add_formatted_text(para, stripped)

            # Empty line
            else:
                if in_table:
                    self._add_table_to_docx(doc, table_data)
                    in_table = False
                    table_data = []

        # Handle remaining table
        if in_table and table_data:
            self._add_table_to_docx(doc, table_data)

    def _add_table_to_docx(self, doc, table_data: List[List[str]]):
        """Add a table to the DOCX document."""
        if not table_data:
            return

        rows = len(table_data)
        cols = max(len(row) for row in table_data)

        table = doc.add_table(rows=rows, cols=cols)
        table.style = 'Table Grid'

        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                if j < cols:
                    table.rows[i].cells[j].text = cell_text

                    # Bold header row
                    if i == 0:
                        for paragraph in table.rows[i].cells[j].paragraphs:
                            for run in paragraph.runs:
                                run.bold = True

        doc.add_paragraph()

    def _add_formatted_text(self, paragraph, text: str):
        """Add text with basic Markdown formatting to a paragraph."""
        import re

        # Simple pattern for bold and italic
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*|`.*?`)', text)

        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = paragraph.add_run(part[2:-2])
                run.bold = True
            elif part.startswith('*') and part.endswith('*') and not part.startswith('**'):
                run = paragraph.add_run(part[1:-1])
                run.italic = True
            elif part.startswith('`') and part.endswith('`'):
                run = paragraph.add_run(part[1:-1])
                run.font.name = 'Consolas'
            else:
                paragraph.add_run(part)

    async def _export_pdf(self, document: GeneratedDocument, organization_name: str = "Organisation") -> tuple[bytes, str]:
        """Export to PDF format using WeasyPrint with PSSIG styling."""
        try:
            from weasyprint import HTML, CSS

            # Get HTML content with organization name
            html_content, _ = await self._export_html(document, organization_name)

            # Convert to PDF
            html = HTML(string=html_content.decode('utf-8'))
            pdf_content = html.write_pdf()

            filename = f"{document.code}_v{document.version}.pdf"
            return pdf_content, filename

        except ImportError:
            # Fallback: return HTML with PDF note
            logger.warning("WeasyPrint not available, returning HTML instead of PDF")
            html_content, _ = await self._export_html(document, organization_name)
            return html_content, f"{document.code}_v{document.version}.html"

    async def _export_xlsx(self, document: GeneratedDocument) -> tuple[bytes, str]:
        """Export to XLSX format (for registers and matrices)."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = Workbook()
        ws = wb.active
        ws.title = document.code[:31]  # Max 31 chars for sheet name

        # Styles
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4FACFE", end_color="4FACFE", fill_type="solid")
        border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )

        # Add metadata
        ws['A1'] = "Document:"
        ws['B1'] = document.code
        ws['A2'] = "Version:"
        ws['B2'] = document.version
        ws['A3'] = "Type:"
        ws['B3'] = document.document_type.value
        ws['A4'] = "Date:"
        ws['B4'] = document.created_at.strftime('%d/%m/%Y')

        # Parse tables from markdown
        tables = self._extract_tables_from_markdown(document.content_markdown or "")

        current_row = 6

        if tables:
            for table_idx, table in enumerate(tables):
                if table:
                    # Add table header
                    for col_idx, header in enumerate(table[0], 1):
                        cell = ws.cell(row=current_row, column=col_idx, value=header)
                        cell.font = header_font
                        cell.fill = header_fill
                        cell.border = border
                        cell.alignment = Alignment(horizontal='center')

                    current_row += 1

                    # Add data rows
                    for row_data in table[1:]:
                        for col_idx, value in enumerate(row_data, 1):
                            cell = ws.cell(row=current_row, column=col_idx, value=value)
                            cell.border = border
                        current_row += 1

                    current_row += 2  # Space between tables
        else:
            # No tables found, create content sheet
            ws.cell(row=6, column=1, value="Contenu du document:")
            lines = (document.content_markdown or "").split('\n')
            for i, line in enumerate(lines[:100], 7):  # Limit to 100 lines
                ws.cell(row=i, column=1, value=line)

        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            ws.column_dimensions[column_letter].width = min(max_length + 2, 50)

        # Save to bytes
        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)

        filename = f"{document.code}_v{document.version}.xlsx"
        return buffer.getvalue(), filename

    async def _export_csv(self, document: GeneratedDocument) -> tuple[bytes, str]:
        """Export to CSV format (for registers)."""
        import csv

        tables = self._extract_tables_from_markdown(document.content_markdown or "")

        buffer = io.StringIO()
        writer = csv.writer(buffer)

        # Write metadata
        writer.writerow(["# Document", document.code])
        writer.writerow(["# Version", document.version])
        writer.writerow(["# Type", document.document_type.value])
        writer.writerow([])

        if tables:
            for table in tables:
                for row in table:
                    writer.writerow(row)
                writer.writerow([])
        else:
            # Write content as single column
            lines = (document.content_markdown or "").split('\n')
            for line in lines:
                writer.writerow([line])

        filename = f"{document.code}_v{document.version}.csv"
        return buffer.getvalue().encode('utf-8'), filename

    async def _export_pptx(self, document: GeneratedDocument) -> tuple[bytes, str]:
        """Export to PPTX format (for presentations and schemas)."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = document.name
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 26, 46)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{document.code} | Version {document.version}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(136, 136, 136)
        p2.alignment = PP_ALIGN.CENTER

        # Content slides
        sections = self._extract_sections_from_markdown(document.content_markdown or "")

        for section_title, section_content in sections[:10]:  # Limit to 10 slides
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Section title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 26, 46)

            # Content
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True

            # Add content lines
            lines = section_content.split('\n')[:15]  # Limit lines per slide
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = line.strip()
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(51, 51, 51)

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        filename = f"{document.code}_v{document.version}.pptx"
        return buffer.getvalue(), filename

    def _extract_tables_from_markdown(self, content: str) -> List[List[List[str]]]:
        """Extract tables from Markdown content."""
        tables = []
        current_table = []
        in_table = False

        for line in content.split('\n'):
            stripped = line.strip()
            if stripped.startswith('|'):
                # Skip separator lines
                if stripped.replace('|', '').replace('-', '').replace(' ', '').replace(':', ''):
                    in_table = True
                    row = [cell.strip() for cell in stripped.split('|')[1:-1]]
                    if row:
                        current_table.append(row)
            else:
                if in_table and current_table:
                    tables.append(current_table)
                    current_table = []
                in_table = False

        if current_table:
            tables.append(current_table)

        return tables

    def _extract_sections_from_markdown(self, content: str) -> List[tuple]:
        """Extract sections (h1/h2) from Markdown content."""
        sections = []
        current_title = "Introduction"
        current_content = []

        for line in content.split('\n'):
            if line.startswith('# '):
                if current_content:
                    sections.append((current_title, '\n'.join(current_content)))
                current_title = line[2:].strip()
                current_content = []
            elif line.startswith('## '):
                if current_content:
                    sections.append((current_title, '\n'.join(current_content)))
                current_title = line[3:].strip()
                current_content = []
            else:
                current_content.append(line)

        if current_content:
            sections.append((current_title, '\n'.join(current_content)))

        return sections

    def _markdown_to_html(self, markdown_content: str) -> str:
        """Convert Markdown to HTML."""
        try:
            import markdown
            return markdown.markdown(
                markdown_content or "",
                extensions=['tables', 'toc', 'fenced_code']
            )
        except ImportError:
            return f"<pre>{markdown_content}</pre>"


# Singleton instance
export_service = ExportService()
